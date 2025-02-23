# LICENSE HEADER MANAGED BY add-license-header
#
# Copyright 2024-2025 Syntropix
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
#

import os
import time
from pathlib import Path
from typing import Optional

import requests
from pptx import Presentation

from synthora.toolkits.base import BaseToolkit
from synthora.toolkits.decorators import tool
from synthora.types.enums import Ok, Result


class SlidesToolkit(BaseToolkit):
    r"""A toolkit for creating slides"""

    def __init__(
        self,
        template_dir: Optional[str] = None,
        cache_dir: Optional[str] = None,
        tmp_name: Optional[str] = None,
        file_name: Optional[str] = None,
        api_key: Optional[str] = None,
    ):
        super().__init__()
        cwd = Path(os.getcwd())
        self.api_key = api_key or os.getenv("UNSPLASH_API_KEY")
        self.file_name = file_name
        self.template_dir = template_dir
        self.cache_dir = cache_dir or cwd / "cache"
        self.image_bed_pattern = f"https://api.unsplash.com/search/photos?page=1&client_id={self.api_key}&query="

        if not os.path.exists(self.cache_dir):
            os.makedirs(self.cache_dir)

        self.ppt_file: Presentation = None

        self.ppt_template_names = []
        all_files = os.listdir(self.template_dir)
        for file_name in all_files:
            if file_name.lower().endswith(".pptx"):
                self.ppt_template_names.append(file_name.split(".")[0])
        self.tmp_ppt = os.path.join(self.cache_dir, tmp_name or "cache.pptx")
        if os.path.exists(self.tmp_ppt):
            self.ppt_file = Presentation(self.tmp_ppt)

    @tool
    def get_template_names(self) -> Result[str, Exception]:
        r"""Get the names of the available pptx templates."""
        return Ok(str(self.ppt_template_names))

    @staticmethod
    def _return_timestamp() -> str:
        return str(time.time())

    @tool
    def create_file(
        self, theme: Optional[str] = None
    ) -> Result[str, Exception]:
        r"""Create a ppt file with the given theme.

        Args:
            theme:
                The theme of the ppt file to create.
        """
        if self.template_dir and theme in self.ppt_template_names:
            self.ppt_file = Presentation(
                os.path.join(self.template_dir, f"{theme}.pptx")
            )
        else:
            self.ppt_file = Presentation()
        self.ppt_file.save(os.path.join(self.cache_dir, self.tmp_ppt))
        return Ok(
            f"created a ppt file with theme: {theme},"
            + f" saved at {os.path.join(self.cache_dir, self.tmp_ppt)}"
        )

    @tool
    def get_image(self, keywords: str) -> Result[str, Exception]:
        r"""Get an image from the internet based on the keywords.

        Args:
            keywords:
                The keywords to search for the image.
                Must in English.
        """
        print(keywords)
        picture_url = self.image_bed_pattern + keywords
        response = requests.get(picture_url).json()
        if not response:
            return Ok("no image found")
        img_url = response["results"][0]["urls"]["regular"]
        image = requests.get(img_url)
        img_local_path = os.path.join(
            self.cache_dir, f"{self._return_timestamp()}.jpg"
        )
        with open(img_local_path, "wb") as f:
            f.write(image.content)
        return Ok(f"image saved at {img_local_path}")

    @tool
    def add_first_page(
        self, title: str, subtitle: str
    ) -> Result[str, Exception]:
        r"""Add the first page to the ppt file.

        Args:
            title:
                The title of the first page.
            subtitle:
                The subtitle of the first page.
        """

        slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[0])
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]

        title_shape.text = title
        subtitle_shape.text = subtitle
        self.ppt_file.save(os.path.join(self.cache_dir, self.tmp_ppt))
        return Ok("added first page")

    @tool
    def add_text_page(
        self, title: str, bullet_items: str
    ) -> Result[str, Exception]:
        r"""Add a text page(outline page is also applied) to the ppt file.

        Args:
            title:
                The title of the text page.
            bullet_items:
                The bullet items of the text page. Should be string,
                for multiple bullet items, please use [SPAN] to separate them.

        """
        slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame

        _bullet_items = bullet_items.split("[SPAN]")
        for bullet_item in _bullet_items:
            bullet_item_strip = bullet_item.strip()
            p = tf.add_paragraph()
            p.text = bullet_item_strip
            p.level = 1
        self.ppt_file.save(os.path.join(self.cache_dir, self.tmp_ppt))
        return Ok("added page")

    @tool
    def add_text_image_page(
        self, title: str, bullet_items: str, image: str
    ) -> Result[str, Exception]:
        r"""Add a text page with one image to the ppt file.

        Args:
            title:
                The title of the text page.
            bullet_items:
                The bullet items of the text page. Should be string,
                for multiple bullet items, please use [SPAN] to separate them.
            image:
                The image to add to the text page. Should be a path.
        """

        slide = self.ppt_file.slides.add_slide(self.ppt_file.slide_layouts[3])
        title_shape = slide.shapes.title
        title_shape.text = title

        body_shape = slide.placeholders[1]

        tf = body_shape.text_frame
        _bullet_items = bullet_items.split("[SPAN]")
        for bullet_item in _bullet_items:
            bullet_item_strip = bullet_item.strip()
            p = tf.add_paragraph()
            p.text = bullet_item_strip
            p.level = 1

        image_shape = slide.placeholders[2]
        slide.shapes.add_picture(
            image,
            image_shape.left,
            image_shape.top,
            image_shape.width,
            image_shape.height,
        )
        self.ppt_file.save(os.path.join(self.cache_dir, self.tmp_ppt))
        return Ok("added page")

    @tool
    def submit_file(self) -> Result[str, Exception]:
        r"""Finalize the slides file and submit it.
        You must use this function to submit your work.
        """

        file_path = os.path.join(
            self.cache_dir,
            f"{self.file_name or self._return_timestamp()}.pptx",
        )
        self.ppt_file.save(file_path)
        os.remove(os.path.join(self.cache_dir, self.tmp_ppt))
        return Ok(f"submitted. view slides at {file_path}")
